from __future__ import annotations

import math
import os
import subprocess
from pathlib import Path
from typing import Dict, Iterable, List, Tuple

import pandas as pd
from pypdf import PdfReader
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Inches, Pt


ROOT = Path(__file__).resolve().parent
PDF_PATH = ROOT / "本科毕业论文  王浩.pdf"
OUTPUT_PATH = ROOT / "毕业答辩PPT_王浩.pptx"

TITLE = "LSSeg 与 SAM-LoRA 级联的医学线状结构图像分割方法研究"
AUTHOR = "王浩"
ADVISOR = "赵爱特"
COLLEGE = "计算机科学技术学院"
MAJOR = "计算机科学与技术"
DATE_TEXT = "2026 年 5 月"

DATASETS = ["STARE", "RITE", "CHASE_DB1", "HRF", "AxonDeepSeg_SEM"]
PAIR_LOGS = {
    "STARE": {
        "baseline": ROOT / "log/test_STARE_LSSeg 02-13 11_34/final_metrics.csv",
        "cascade": ROOT / "log/test_STARE_LSSegSAMLoRA_DiceCE 04-10 23_07/final_metrics.csv",
    },
    "RITE": {
        "baseline": ROOT / "log/test_RITE_LSSeg 02-14 08_30/final_metrics.csv",
        "cascade": ROOT / "log/test_RITE_LSSegSAMLoRA 04-03 14_33/final_metrics.csv",
    },
    "CHASE_DB1": {
        "baseline": ROOT / "log/test_CHASE_DB1_LSSeg 02-17 12_48/final_metrics.csv",
        "cascade": ROOT / "log/test_CHASE_DB1_LSSegSAMLoRA_DiceCE 04-12 00_09/final_metrics.csv",
    },
    "HRF": {
        "baseline": ROOT / "log/test_HRF_LSSeg 02-14 21_39/final_metrics.csv",
        "cascade": ROOT / "log/test_HRF_LSSegSAMLoRA_DiceCE 04-05 21_24/final_metrics.csv",
    },
    "AxonDeepSeg_SEM": {
        "baseline": ROOT / "log/test_AxonDeepSeg_SEM_LSSeg 03-02 15_55/final_metrics.csv",
        "cascade": ROOT / "log/test_AxonDeepSeg_SEM_LSSegSAMLoRA_DiceCE 04-18 09_22/final_metrics.csv",
    },
}

FIGURES = {
    "framework": ROOT / "figures/thesis-research-framework.svg",
    "architecture": ROOT / "figures/lssseg-sam-lora-architecture.svg",
    "dataflow": ROOT / "figures/lssseg-sam-lora-dataflow-v2.svg",
    "lsseg": ROOT / "figures/lsseg-principle.svg",
    "sam": ROOT / "figures/sam-principle.svg",
    "lora": ROOT / "figures/lora-principle.svg",
    "fusion": ROOT / "figures/residual-fusion-conv-principle.svg",
    "metric_compare": ROOT / "figures/selected_log_metrics_comparison.png",
    "vignettes": [
        ROOT / "figures/best_worst_showcase/STARE_best_worst_showcase.png",
        ROOT / "figures/best_worst_showcase/RITE_best_worst_showcase.png",
        ROOT / "figures/best_worst_showcase/CHASE_DB1_best_worst_showcase.png",
        ROOT / "figures/best_worst_showcase/HRF_best_worst_showcase.png",
        ROOT / "figures/best_worst_showcase/AxonDeepSeg_SEM_best_worst_showcase.png",
    ],
}


def read_pdf_title(pdf_path: Path) -> str:
    reader = PdfReader(str(pdf_path))
    first = reader.pages[1].extract_text() or ""
    for line in first.splitlines():
        line = line.strip()
        if line and "LSSeg" in line and "SAM-LoRA" in line:
            return line.replace("  ", " ").strip()
    return TITLE


def read_metrics(csv_path: Path) -> Dict[str, float]:
    df = pd.read_csv(csv_path)
    if "Unnamed: 0" in df.columns:
        key_col = "Unnamed: 0"
    else:
        key_col = df.columns[0]
    metrics = {}
    for _, row in df.iterrows():
        metrics[str(row[key_col])] = float(row["Mean"])
    return metrics


def build_metric_table() -> List[Dict[str, float]]:
    rows = []
    for dataset in DATASETS:
        base = read_metrics(PAIR_LOGS[dataset]["baseline"])
        casc = read_metrics(PAIR_LOGS[dataset]["cascade"])
        rows.append(
            {
                "dataset": dataset,
                "baseline_dsc": base["DSC"],
                "cascade_dsc": casc["DSC"],
                "delta_dsc": casc["DSC"] - base["DSC"],
                "baseline_hd95": base["95HD"],
                "cascade_hd95": casc["95HD"],
                "delta_hd95": casc["95HD"] - base["95HD"],
            }
        )
    return rows


def ensure_png_from_svg(svg_path: Path, png_path: Path) -> Path:
    if png_path.exists() and png_path.stat().st_mtime >= svg_path.stat().st_mtime:
        return png_path

    png_path.parent.mkdir(parents=True, exist_ok=True)
    subprocess.run(
        [
            "libreoffice",
            "--headless",
            "--convert-to",
            "png",
            "--outdir",
            str(png_path.parent),
            str(svg_path),
        ],
        check=True,
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
    )
    return png_path


def ppt_image_path(path: Path) -> Path:
    if path.suffix.lower() != ".svg":
        return path
    candidate = ROOT / "figures/ppt_exports" / f"{path.stem}.png"
    if candidate.exists():
        return candidate
    fallback = path.with_suffix(".png")
    if fallback.exists():
        return fallback
    return path


def set_background(slide, color: str = "F7FAFC"):
    fill = slide.background.fill
    fill.solid()
    from pptx.dml.color import RGBColor

    fill.fore_color.rgb = RGBColor.from_string(color)


def add_header(slide, title: str, subtitle: str | None = None):
    box = slide.shapes.add_textbox(Cm(1.0), Cm(0.7), Cm(31.0), Cm(1.4))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(24)
    r.font.bold = True
    from pptx.dml.color import RGBColor

    r.font.color.rgb = RGBColor.from_string("102A43")
    p.alignment = PP_ALIGN.LEFT
    if subtitle:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(2)
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = "Microsoft YaHei"
        r2.font.size = Pt(11)
        from pptx.dml.color import RGBColor

        r2.font.color.rgb = RGBColor.from_string("627D98")


def add_textbox(slide, left, top, width, height, text="", font_size=18, bold=False, color="1F2937",
                fill=None, line_color="D8E1E8", radius=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    shape_type = SHAPE.ROUNDED_RECTANGLE if radius else SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill is None:
        shp.fill.background()
    else:
        from pptx.dml.color import RGBColor

        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor.from_string(fill)
    from pptx.dml.color import RGBColor

    shp.line.color.rgb = RGBColor.from_string(line_color)
    tf = shp.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    if text:
        r = p.add_run()
        r.text = text
        r.font.name = "Microsoft YaHei"
        r.font.size = Pt(font_size)
        r.font.bold = bold
        r.font.color.rgb = RGBColor.from_string(color)
    p.alignment = align
    return shp


def add_bullets(slide, left, top, width, height, bullets, font_size=16, color="1F2937", title=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    if title:
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.name = "Microsoft YaHei"
        r.font.size = Pt(font_size + 1)
        r.font.bold = True
        from pptx.dml.color import RGBColor

        r.font.color.rgb = RGBColor.from_string(color)
        p.space_after = Pt(6)
        start = 1
    else:
        start = 0
    from pptx.dml.color import RGBColor

    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[start] if start == 0 and idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor.from_string(color)
        p.space_after = Pt(4)
    return box


def add_image(slide, path: Path, left, top, width=None, height=None):
    path = ppt_image_path(path)
    if not path.exists():
        return None
    if path.suffix.lower() == ".svg":
        return None
    if width is not None and height is not None:
        return slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    if width is not None:
        return slide.shapes.add_picture(str(path), left, top, width=width)
    if height is not None:
        return slide.shapes.add_picture(str(path), left, top, height=height)
    return slide.shapes.add_picture(str(path), left, top)


def add_footer(slide, text="王浩 | 毕业答辩"):
    box = slide.shapes.add_textbox(Cm(1.0), Cm(18.3), Cm(31.0), Cm(0.5))
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.runs[0]
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(10)
    from pptx.dml.color import RGBColor

    run.font.color.rgb = RGBColor.from_string("667085")


def add_title_slide(prs, title: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, "F7FAFC")

    add_textbox(
        slide,
        Cm(1.2),
        Cm(1.2),
        Cm(17.5),
        Cm(4.0),
        text=title,
        font_size=28,
        bold=True,
        color="102A43",
        fill="FFFFFF",
        line_color="E2E8F0",
        radius=True,
    )
    add_textbox(
        slide,
        Cm(1.55),
        Cm(2.35),
        Cm(15.5),
        Cm(1.6),
        text="基于 LSSeg 与 SAM-LoRA 级联的生物医学线状结构分割方法研究",
        font_size=20,
        bold=False,
        color="334E68",
    )
    add_textbox(
        slide,
        Cm(1.55),
        Cm(4.0),
        Cm(15.8),
        Cm(0.9),
        text=f"{AUTHOR}    {COLLEGE}    {MAJOR}",
        font_size=14,
        color="486581",
    )
    add_textbox(
        slide,
        Cm(1.55),
        Cm(4.8),
        Cm(15.8),
        Cm(0.9),
        text=f"指导教师：{ADVISOR}    答辩时间：{DATE_TEXT}",
        font_size=14,
        color="486581",
    )

    add_textbox(
        slide,
        Cm(19.8),
        Cm(1.4),
        Cm(11.5),
        Cm(10.0),
        fill="E6FFFB",
        line_color="B2F5EA",
        radius=True,
    )
    add_textbox(
        slide,
        Cm(20.5),
        Cm(2.0),
        Cm(10.1),
        Cm(1.0),
        text="答辩主线",
        font_size=22,
        bold=True,
        color="0F766E",
        fill="E6FFFB",
        line_color="E6FFFB",
    )
    bullets = [
        "研究对象是血管、轴突等线状结构图像",
        "核心方法是 LSSeg 生成粗分割，SAM-LoRA 完成细化修正",
        "实验覆盖 STARE、RITE、CHASE_DB1、HRF、AxonDeepSeg_SEM 五个数据集",
    ]
    add_bullets(slide, Cm(20.3), Cm(3.1), Cm(10.6), Cm(3.2), bullets, font_size=14, color="134E4A")
    add_footer(slide)
    return slide


def add_simple_slide(prs, title, bullets, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_textbox(slide, Cm(0.8), Cm(0.4), Cm(31.2), Cm(1.0), text=title, font_size=24, bold=True, color="102A43", fill=None, line_color="FFFFFF")
    if note:
        add_textbox(slide, Cm(0.85), Cm(1.35), Cm(30.2), Cm(0.7), text=note, font_size=11, color="627D98", fill=None, line_color="FFFFFF")
    add_bullets(slide, Cm(1.0), Cm(2.0), Cm(29.8), Cm(13.5), bullets, font_size=18)
    add_footer(slide)
    return slide


def add_pipeline_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_textbox(slide, Cm(0.8), Cm(0.4), Cm(31.2), Cm(1.0), text="研究思路与技术路线", font_size=24, bold=True, color="102A43", fill=None, line_color="FFFFFF")
    add_textbox(slide, Cm(0.9), Cm(1.3), Cm(30.0), Cm(0.8), text="这页只保留主线：粗分割、提示构造、细化修正、残差融合、实验验证。", font_size=11, color="627D98", fill=None, line_color="FFFFFF")

    xs = [1.0, 6.0, 11.0, 16.0, 21.0, 26.0]
    labels = [
        ("数据集", "STARE / RITE / CHASE_DB1 / HRF / AxonDeepSeg_SEM"),
        ("LSSeg", "生成稳定粗分割"),
        ("Prompt", "Mask / Box prompt"),
        ("SAM-LoRA", "精细修正"),
        ("Fusion", "残差融合"),
        ("输出", "分割结果与指标"),
    ]
    for idx, (x, (title, desc)) in enumerate(zip(xs, labels)):
        add_textbox(
            slide,
            Cm(x),
            Cm(3.0),
            Cm(4.0),
            Cm(2.2),
            text=title,
            font_size=18,
            bold=True,
            color="0F172A",
            fill="FFFFFF",
            line_color="D9E2EC",
            radius=True,
        )
        add_textbox(
            slide,
            Cm(x + 0.15),
            Cm(3.7),
            Cm(3.7),
            Cm(1.1),
            text=desc,
            font_size=11,
            color="486581",
            fill="FFFFFF",
            line_color="FFFFFF",
            radius=True,
        )
        if idx < len(xs) - 1:
            arrow = slide.shapes.add_shape(SHAPE.RIGHT_ARROW, Cm(x + 4.0), Cm(3.75), Cm(0.8), Cm(0.7))
            arrow.fill.solid()
            from pptx.dml.color import RGBColor

            arrow.fill.fore_color.rgb = RGBColor.from_string("94A3B8")
            arrow.line.color.rgb = RGBColor.from_string("94A3B8")

    # simple arrows using lines
    for x in [5.1, 10.1, 15.1, 20.1, 25.1]:
        line = slide.shapes.add_shape(SHAPE.RIGHT_ARROW, Cm(x), Cm(3.65), Cm(0.8), Cm(0.9))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor.from_string("94A3B8")
        line.line.color.rgb = RGBColor.from_string("94A3B8")

    add_textbox(slide, Cm(0.9), Cm(6.1), Cm(14.2), Cm(6.0), fill="FFFFFF", line_color="D9E2EC", radius=True)
    add_textbox(slide, Cm(1.2), Cm(6.45), Cm(13.6), Cm(0.6), text="研究重点", font_size=18, bold=True, color="0F766E")
    add_bullets(
        slide,
        Cm(1.2),
        Cm(7.0),
        Cm(13.2),
        Cm(4.4),
        [
            "围绕生物医学线状结构分割任务展开。",
            "把 LSSeg 和 SAM-LoRA 组合成级联框架。",
            "重点解决 prompt 构造、训练稳定性和边界修正问题。",
            "使用多数据集交叉验证与结构化指标评估方法有效性。",
        ],
        font_size=15,
        color="334E68",
    )
    add_textbox(slide, Cm(15.5), Cm(6.1), Cm(15.0), Cm(6.0), fill="F8FAFC", line_color="D9E2EC", radius=True)
    add_textbox(slide, Cm(15.8), Cm(6.45), Cm(14.2), Cm(0.6), text="可复现的技术栈", font_size=18, bold=True, color="0F766E")
    add_bullets(
        slide,
        Cm(15.8),
        Cm(7.0),
        Cm(13.9),
        Cm(4.4),
        [
            "PyTorch + MONAI + Optuna",
            "Albumentations 数据增强",
            "DiceCE 损失、AdamW 优化器、StepLR 调度",
            "AMP 混合精度训练与 5 折交叉验证",
        ],
        font_size=15,
        color="334E68",
    )
    add_footer(slide)
    return slide


def add_diagram_slide(prs, title, left_caption, image_path: Path, right_caption, text_blocks):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_textbox(slide, Cm(0.8), Cm(0.4), Cm(31.2), Cm(1.0), text=title, font_size=24, bold=True, color="102A43", fill=None, line_color="FFFFFF")
    add_image(slide, image_path, Cm(1.0), Cm(1.7), width=Cm(14.2))
    add_textbox(slide, Cm(1.0), Cm(12.2), Cm(14.2), Cm(0.6), text=left_caption, font_size=11, color="64748B", fill=None, line_color="FFFFFF", align=PP_ALIGN.CENTER)
    add_textbox(slide, Cm(16.0), Cm(1.7), Cm(14.0), Cm(10.8), fill="FFFFFF", line_color="D9E2EC", radius=True)
    add_textbox(slide, Cm(16.35), Cm(2.0), Cm(13.1), Cm(0.7), text=right_caption, font_size=18, bold=True, color="0F766E")
    add_bullets(slide, Cm(16.35), Cm(2.8), Cm(13.0), Cm(8.5), text_blocks, font_size=15, color="334E68")
    add_footer(slide)
    return slide


def add_datasets_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_textbox(slide, Cm(0.8), Cm(0.4), Cm(31.2), Cm(1.0), text="实验数据集与评价指标", font_size=24, bold=True, color="102A43", fill=None, line_color="FFFFFF")
    add_textbox(slide, Cm(0.9), Cm(1.25), Cm(29.8), Cm(0.7), text="项目最终答辩只使用这 5 个数据集，和论文中的对比实验保持一致。", font_size=11, color="627D98", fill=None, line_color="FFFFFF")

    table_data = [
        ("STARE", "20", "视网膜血管"),
        ("RITE", "40", "视网膜血管"),
        ("CHASE_DB1", "28", "视网膜血管"),
        ("HRF", "45", "高分辨率视网膜血管"),
        ("AxonDeepSeg_SEM", "10", "扫描电镜轴突"),
    ]
    rows = len(table_data) + 1
    cols = 3
    table = slide.shapes.add_table(rows, cols, Cm(1.0), Cm(2.0), Cm(13.5), Cm(8.0)).table
    headers = ["数据集", "样本数", "任务类型"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
    for i, row in enumerate(table_data, start=1):
        for j, val in enumerate(row):
            table.cell(i, j).text = val
    for j, w in enumerate([4.4, 2.8, 6.3]):
        table.columns[j].width = Cm(w)
    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i, j)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.name = "Microsoft YaHei"
                    r.font.size = Pt(12 if i else 13)
                    r.font.bold = i == 0
                    if i == 0:
                        from pptx.dml.color import RGBColor
                        r.font.color.rgb = RGBColor.from_string("0F766E")
            if i == 0:
                cell.fill.solid()
                from pptx.dml.color import RGBColor
                cell.fill.fore_color.rgb = RGBColor.from_string("D9F2EE")

    add_textbox(slide, Cm(15.0), Cm(2.0), Cm(15.0), Cm(3.0), fill="FFFFFF", line_color="D9E2EC", radius=True)
    add_textbox(slide, Cm(15.35), Cm(2.3), Cm(13.8), Cm(0.5), text="评价指标", font_size=18, bold=True, color="0F766E")
    add_bullets(
        slide,
        Cm(15.35),
        Cm(3.0),
        Cm(13.3),
        Cm(1.8),
        ["AP", "DSC", "mIoU", "clDice", "ASSD", "95HD", "ODS"],
        font_size=14,
        color="334E68",
    )
    add_textbox(slide, Cm(15.0), Cm(5.4), Cm(15.0), Cm(4.6), fill="FFFFFF", line_color="D9E2EC", radius=True)
    add_textbox(slide, Cm(15.35), Cm(5.75), Cm(13.8), Cm(0.5), text="评估说明", font_size=18, bold=True, color="0F766E")
    add_bullets(
        slide,
        Cm(15.35),
        Cm(6.4),
        Cm(13.4),
        Cm(2.8),
        [
            "区域指标看分割是否覆盖目标。",
            "clDice 看细长结构是否连通。",
            "ASSD 与 95HD 看边界是否贴合。",
            "ODS 用统一阈值评价整体检测效果。",
        ],
        font_size=14,
        color="334E68",
    )
    add_footer(slide)
    return slide


def add_quant_slide(prs, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_textbox(slide, Cm(0.8), Cm(0.4), Cm(31.2), Cm(1.0), text="定量实验结果", font_size=24, bold=True, color="102A43", fill=None, line_color="FFFFFF")
    add_textbox(slide, Cm(0.9), Cm(1.25), Cm(29.8), Cm(0.7), text="图表来自项目中已经选定的实验日志，和论文第四章保持一致。", font_size=11, color="627D98", fill=None, line_color="FFFFFF")

    add_image(slide, FIGURES["metric_compare"], Cm(1.0), Cm(1.8), width=Cm(17.2))

    table_rows = len(rows) + 1
    table = slide.shapes.add_table(table_rows, 4, Cm(18.5), Cm(1.9), Cm(11.6), Cm(10.1)).table
    headers = ["数据集", "DSC↑", "95HD↓", "结论"]
    for j, h in enumerate(headers):
        table.cell(0, j).text = h
    for i, row in enumerate(rows, start=1):
        table.cell(i, 0).text = row["dataset"]
        table.cell(i, 1).text = f'{row["delta_dsc"]:+.3f}'
        table.cell(i, 2).text = f'{row["delta_hd95"]:+.3f}'
        table.cell(i, 3).text = "改进" if row["delta_dsc"] >= 0 else "下降"
    for j, w in enumerate([3.0, 2.0, 2.2, 4.4]):
        table.columns[j].width = Cm(w)
    for i in range(table_rows):
        for j in range(4):
            cell = table.cell(i, j)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.name = "Microsoft YaHei"
                    r.font.size = Pt(11 if i else 12)
                    r.font.bold = i == 0
                    if i == 0:
                        from pptx.dml.color import RGBColor
                        r.font.color.rgb = RGBColor.from_string("0F766E")
            if i == 0:
                cell.fill.solid()
                from pptx.dml.color import RGBColor
                cell.fill.fore_color.rgb = RGBColor.from_string("D9F2EE")

    add_footer(slide)
    return slide


def add_qualitative_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_textbox(slide, Cm(0.8), Cm(0.4), Cm(31.2), Cm(1.0), text="定性结果展示", font_size=24, bold=True, color="102A43", fill=None, line_color="FFFFFF")
    add_textbox(slide, Cm(0.9), Cm(1.25), Cm(29.8), Cm(0.7), text="每张图是一个数据集的代表性样本，展示 baseline 与 cascade 的对比效果。", font_size=11, color="627D98", fill=None, line_color="FFFFFF")

    positions = [
        (1.0, 2.0),
        (11.0, 2.0),
        (21.0, 2.0),
        (5.5, 9.0),
        (15.5, 9.0),
    ]
    captions = DATASETS
    for (x, y), path, caption in zip(positions, FIGURES["vignettes"], captions):
        add_textbox(slide, Cm(x), Cm(y - 0.4), Cm(8.2), Cm(0.4), text=caption, font_size=13, bold=True, color="0F766E")
        add_image(slide, path, Cm(x), Cm(y), width=Cm(8.2))

    add_textbox(slide, Cm(1.0), Cm(15.8), Cm(29.0), Cm(1.3), text="这些样本里，级联模型主要在细小分支补全、边界贴合和局部连通性上更稳。", font_size=14, color="334E68", fill="FFFFFF", line_color="D9E2EC", radius=True)
    add_footer(slide)
    return slide


def add_conclusion_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_textbox(slide, Cm(0.8), Cm(0.4), Cm(31.2), Cm(1.0), text="结论与展望", font_size=24, bold=True, color="102A43", fill=None, line_color="FFFFFF")
    add_textbox(slide, Cm(1.0), Cm(1.8), Cm(14.3), Cm(9.3), fill="FFFFFF", line_color="D9E2EC", radius=True)
    add_textbox(slide, Cm(1.35), Cm(2.1), Cm(13.4), Cm(0.6), text="结论", font_size=18, bold=True, color="0F766E")
    add_bullets(
        slide,
        Cm(1.35),
        Cm(2.8),
        Cm(13.2),
        Cm(6.9),
        [
            "LSSegSAMLoRA 在 5 个数据集上都优于单独 LSSeg。",
            "改进主要体现在细小结构补全、边界修正和连通性恢复。",
            "残差融合和 prompt 优化让通用分割模型更好适配医学图像。",
            "方法适合生物医学线状结构分割这一类任务。",
        ],
        font_size=15,
        color="334E68",
    )
    add_textbox(slide, Cm(16.0), Cm(1.8), Cm(14.0), Cm(9.3), fill="F8FAFC", line_color="D9E2EC", radius=True)
    add_textbox(slide, Cm(16.35), Cm(2.1), Cm(13.3), Cm(0.6), text="展望", font_size=18, bold=True, color="0F766E")
    add_bullets(
        slide,
        Cm(16.35),
        Cm(2.8),
        Cm(13.0),
        Cm(6.9),
        [
            "继续优化自动 prompt 构造方式。",
            "增强融合模块的自适应能力。",
            "扩展到更多医学图像模态和更多场景。",
            "补充更系统的消融实验和推理效率分析。",
        ],
        font_size=15,
        color="334E68",
    )
    add_footer(slide)
    return slide


def add_thanks_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, "F7FAFC")
    add_textbox(
        slide,
        Cm(3.0),
        Cm(4.0),
        Cm(24.0),
        Cm(2.5),
        text="谢谢老师们的指导",
        font_size=30,
        bold=True,
        color="102A43",
        fill="FFFFFF",
        line_color="E2E8F0",
        radius=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_textbox(
        slide,
        Cm(4.0),
        Cm(6.2),
        Cm(22.0),
        Cm(0.8),
        text=f"{AUTHOR} | {TITLE}",
        font_size=14,
        color="627D98",
        fill=None,
        line_color="FFFFFF",
        align=PP_ALIGN.CENTER,
    )
    return slide


def main():
    title = read_pdf_title(PDF_PATH)
    rows = build_metric_table()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, title)
    add_simple_slide(
        prs,
        "研究背景",
        [
            "血管、神经纤维、轴突等线状结构在医学分析中很重要，但通常细小、分支多、边界模糊。",
            "现有分割方法在细小结构保持、局部连通性恢复和复杂边界刻画方面还有不足。",
            "通用分割基础模型具备较强泛化能力，但直接迁移到医学图像时通常需要任务适配。",
        ],
        note="这一页对应论文绪论和摘要中的问题定义。",
    )
    add_simple_slide(
        prs,
        "研究内容",
        [
            "构建 LSSeg 与 SAM-LoRA 的级联分割框架。",
            "研究自动 prompt 构造策略，包括 mask prompt、box prompt 和 prompt_bias。",
            "设计残差融合机制，把粗分割结果与细化结果整合起来。",
            "在 5 个公开数据集上验证该方法的有效性。",
        ],
        note="这一页可以直接作为答辩时的研究内容页。",
    )
    add_pipeline_slide(prs)
    add_diagram_slide(
        prs,
        "核心方法一：LSSeg 粗分割",
        "LSSeg 原理图",
        FIGURES["lsseg"],
        "作用",
        [
            "先从输入图像中提取结构稳定的粗分割结果。",
            "作为后端 SAM-LoRA 的先验提示来源。",
            "在级联框架里承担“保底”和“提纲挈领”的角色。",
        ],
    )
    add_diagram_slide(
        prs,
        "核心方法二：SAM-LoRA 细化分支",
        "SAM 与 LoRA 原理图",
        FIGURES["sam"],
        "关键点",
        [
            "SAM 提供通用分割能力，LoRA 负责参数高效适配。",
            "通过 prompt 把前端粗分割结果送入后端。",
            "在边界、细分支和弱响应区域做局部修正。",
        ],
    )
    add_diagram_slide(
        prs,
        "提示构造与融合策略",
        "残差融合原理图",
        FIGURES["fusion"],
        "论文里的两个关键设计",
        [
            "prompt_bias 提高召回，让更多疑似结构被送入 SAM。",
            "box prompt 和 mask prompt 协同使用，增强提示稳定性。",
            "残差融合让最终结果保持 LSSeg 的稳定性，同时吸收 SAM 的修正。",
        ],
    )
    add_datasets_slide(prs)
    add_quant_slide(prs, rows)
    add_qualitative_slide(prs)
    add_conclusion_slide(prs)
    add_thanks_slide(prs)

    prs.save(str(OUTPUT_PATH))
    print(f"saved: {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
